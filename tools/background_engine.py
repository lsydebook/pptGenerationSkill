"""Pure background engine — no text, no content, just visual backgrounds.

Simulates gradients by stacking thin rectangles (python-pptx has no native gradient).
Each layout type gets its own decorative composition.

Standalone test:
    python tools/background_engine.py
"""
from __future__ import annotations

import math
from dataclasses import dataclass, field
from typing import Literal

SlideRole = Literal["cover", "section", "content", "thanks"]


# --- color palette ---
ACCENT_DEEP = "152E4A"
ACCENT_DARK = "1B3A5C"
ACCENT_MAIN = "2B5B84"
ACCENT_MID = "3A7BBF"
ACCENT_LIGHT = "5A9BD5"
BACKGROUND_WARM = "F8FAFE"
BACKGROUND_COOL = "EAF0F8"
PAPER_WHITE = "FFFFFF"
SECTION_BG = "EFF3F9"


SLIDE_W = 12192000
SLIDE_H = 6858000


@dataclass
class BgShape:
    shape_type: str  # "rect" | "oval"
    x: int
    y: int
    w: int
    h: int
    color: str


@dataclass
class BackgroundSpec:
    bg_color: str
    shapes: list[BgShape] = field(default_factory=list)


# ============================================================
# gradient simulation: stack N thin rectangles
# ============================================================

def _hex_to_rgb_ints(h: str) -> tuple[int, int, int]:
    h = h.lstrip("#")
    return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)


def _rgb_to_hex(r: int, g: int, b: int) -> str:
    r = max(0, min(255, int(r)))
    g = max(0, min(255, int(g)))
    b = max(0, min(255, int(b)))
    return f"{r:02X}{g:02X}{b:02X}"


def _gradient_strips(
    color_from: str,
    color_to: str,
    x: int, y: int, w: int, h: int,
    steps: int = 20,
    vertical: bool = True,
) -> list[BgShape]:
    r1, g1, b1 = _hex_to_rgb_ints(color_from)
    r2, g2, b2 = _hex_to_rgb_ints(color_to)
    strips: list[BgShape] = []
    if steps <= 1:
        strips.append(BgShape("rect", x, y, w, h, color_from))
        return strips
    for i in range(steps):
        t = i / (steps - 1)
        cr = r1 + (r2 - r1) * t
        cg = g1 + (g2 - g1) * t
        cb = b1 + (b2 - b1) * t
        color = _rgb_to_hex(cr, cg, cb)
        if vertical:
            sy = y + int(h * i / steps)
            sh = int(h / steps) + 1
            strips.append(BgShape("rect", x, sy, w, sh, color))
        else:
            sx = x + int(w * i / steps)
            sw = int(w / steps) + 1
            strips.append(BgShape("rect", sx, y, sw, h, color))
    return strips


# ============================================================
# layouts
# ============================================================

def _build_gradient_bg_strips() -> list[BgShape]:
    """Full-page subtle gradient: white bottom-left → very pale blue top-right."""
    shapes: list[BgShape] = []
    shapes += _gradient_strips(
        PAPER_WHITE, "EDF2F9",
        x=0, y=0, w=SLIDE_W, h=SLIDE_H,
        steps=24,
        vertical=True,
    )
    return shapes


def cover_background(accent: str = ACCENT_MAIN) -> BackgroundSpec:
    shapes: list[BgShape] = []

    # full solid deep blue
    shapes += _gradient_strips(
        ACCENT_DEEP, ACCENT_DARK,
        x=0, y=0, w=SLIDE_W, h=SLIDE_H,
        steps=1,
        vertical=True,
    )

    # large decorative circles on the right
    circles = [
        (SLIDE_W + 600000, -SLIDE_H // 2, SLIDE_H * 2, SLIDE_H * 2, ACCENT_DARK),
        (SLIDE_W, -800000, SLIDE_H, SLIDE_H, "203860"),
        (SLIDE_W - 1600000, SLIDE_H // 2 - 1200000, 3600000, 3600000, "1E3458"),
        (SLIDE_W - 4000000, SLIDE_H // 4, 2400000, 2400000, "243C62"),
    ]
    for cx, cy, cw, ch, ccol in circles:
        shapes.append(BgShape("oval", int(cx - cw // 2), int(cy - cw // 2), cw, cw, ccol))

    # bottom decorative line
    shapes.append(BgShape("rect", 0, SLIDE_H - 120000, SLIDE_W, 120000, "FFFFFF0D"))

    # thin accent line near center-bottom
    shapes.append(BgShape("rect", SLIDE_W // 2 - 2400000, SLIDE_H - 3600000, 4800000, 30000, "FFFFFF10"))

    return BackgroundSpec(bg_color=ACCENT_DEEP, shapes=shapes)


def section_background(accent: str = ACCENT_MAIN) -> BackgroundSpec:
    shapes: list[BgShape] = []

    # subtle gradient from cool to lighter
    shapes += _gradient_strips(
        "E8EEF7", "F4F7FC",
        x=0, y=0, w=SLIDE_W, h=SLIDE_H,
        steps=20,
        vertical=True,
    )

    # large semi-transparent circles in bottom-right
    circles = [
        (SLIDE_W - 1800000, SLIDE_H + 600000, 4800000, "D9E2F0"),
        (SLIDE_W - 2800000, SLIDE_H + 1800000, 3600000, "D0DBEC"),
    ]
    for cx, cy, cr, ccol in circles:
        shapes.append(BgShape("oval", int(cx - cr // 2), int(cy - cr // 2), cr, cr, ccol))

    # overlapping medium circle
    shapes.append(BgShape("oval", SLIDE_W - 5400000, SLIDE_H - 2400000, 3000000, 3000000, "DDE5F2"))

    return BackgroundSpec(bg_color="F4F7FC", shapes=shapes)


def content_background(accent: str = ACCENT_MAIN) -> BackgroundSpec:
    shapes: list[BgShape] = []

    shapes += _gradient_strips(
        PAPER_WHITE, "F4F7FB",
        x=0, y=0, w=SLIDE_W, h=SLIDE_H,
        steps=16,
        vertical=True,
    )

    return BackgroundSpec(bg_color=PAPER_WHITE, shapes=shapes)


def thanks_background(accent: str = ACCENT_MAIN) -> BackgroundSpec:
    shapes: list[BgShape] = []

    # full solid deep blue
    shapes += _gradient_strips(
        ACCENT_DEEP, ACCENT_DARK,
        x=0, y=0, w=SLIDE_W, h=SLIDE_H,
        steps=1,
        vertical=True,
    )

    # decorative circles
    circles = [
        (SLIDE_W + 800000, -SLIDE_H // 3, SLIDE_H * 2, SLIDE_H * 2, ACCENT_DEEP),
        (SLIDE_W - 1200000, SLIDE_H // 2 - 600000, 3600000, 3600000, "1E3458"),
        (SLIDE_W - 4200000, -1800000, 4800000, 4800000, "203860"),
        (SLIDE_W - 5400000, SLIDE_H - 3000000, 3000000, 3000000, "243C62"),
    ]
    for cx, cy, cw, ch, ccol in circles:
        shapes.append(BgShape("oval", int(cx - cw // 2), int(cy - cw // 2), cw, cw, ccol))

    # horizontal accent lines
    for offset, thickness in [(SLIDE_H // 2 + 1800000, 30000), (SLIDE_H - 240000, 120000)]:
        shapes.append(BgShape("rect", SLIDE_W // 2 - 3600000, offset, 7200000, thickness, "FFFFFF0D"))

    return BackgroundSpec(bg_color=ACCENT_DEEP, shapes=shapes)


# ============================================================
# dispatcher
# ============================================================

def get_background(role: SlideRole, accent: str | None = None) -> BackgroundSpec:
    acc = accent or ACCENT_MAIN
    dispatch = {
        "cover": cover_background,
        "section": section_background,
        "content": content_background,
        "thanks": thanks_background,
    }
    return dispatch.get(role, content_background)(acc)


# ============================================================
# renderer (standalone)
# ============================================================

def render_background_slides(
    output_path: str,
    slide_roles: list[SlideRole],
    accent: str | None = None,
):
    """Render a PPTX file with only backgrounds, no content."""
    from pathlib import Path

    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Emu

    def _hexrgb(h: str) -> RGBColor:
        hh = h.lstrip("#")
        return RGBColor(int(hh[0:2], 16), int(hh[2:4], 16), int(hh[4:6], 16))

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)

    for role in slide_roles:
        bg = get_background(role, accent)
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        if bg.bg_color:
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = _hexrgb(bg.bg_color)

        for s in bg.shapes:
            st = 2 if s.shape_type == "oval" else 1
            shape = slide.shapes.add_shape(st, Emu(s.x), Emu(s.y), Emu(s.w), Emu(s.h))
            shape.fill.solid()
            shape.fill.fore_color.rgb = _hexrgb(s.color)
            shape.line.fill.background()

    prs.save(str(out))
    print(f"saved: {out}")


# ============================================================
# standalone demo
# ============================================================

if __name__ == "__main__":
    roles: list[SlideRole] = [
        "cover",
        "section",
        "content",
        "content",
        "section",
        "content",
        "content",
        "section",
        "content",
        "content",
        "thanks",
    ]
    render_background_slides("out/_bg_demo.pptx", roles)
    print("Done — open out/_bg_demo.pptx to preview backgrounds.")
