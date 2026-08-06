"""Render a list of SlideSpec into a .pptx file using the user's template."""
from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from pptx.presentation import Presentation as PresentationT
from pptx.util import Pt

from schemas.input import InputBundle

from .image_utils import resolve_image_path
from .polish import HighlightStyle, fill_footer, hex_to_rgb, highlight_paragraph
from .pptx_template import load_template, resolve_layout_map


def _strip_existing_slides(prs: PresentationT) -> int:
    sldIdLst = prs.slides._sldIdLst  # noqa: SLF001
    items = list(sldIdLst)
    removed = 0
    for sld_id in items:
        rId = sld_id.rId
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass
        sldIdLst.remove(sld_id)
        removed += 1
    return removed


def _set_text(shape, text, *, highlight=None, font_name=None):
    if text is None:
        return
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.text = text
    for p in tf.paragraphs:
        for run in p.runs:
            if font_name:
                run.font.name = font_name
        if highlight is not None:
            highlight_paragraph(p, highlight)


def _set_bullets(shape, bullets, max_size=28, min_size=14, *, highlight=None, font_name=None):
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.clear()
    bullets = list(bullets)
    if not bullets:
        tf.text = ""
        return

    box_height_emu = shape.height or 0
    if box_height_emu > 0 and len(bullets) > 1:
        box_height_pt = box_height_emu / 12700
        line_height_ratio = 1.5
        space_factor = 1.15
        calc_size = box_height_pt / (len(bullets) * line_height_ratio * space_factor)
        font_size = max(min_size, min(max_size, calc_size))
    else:
        font_size = max_size

    p = tf.paragraphs[0]
    p.text = bullets[0]
    p.level = 0
    for run in p.runs:
        run.font.size = Pt(font_size)
        if font_name:
            run.font.name = font_name
    for item in bullets[1:]:
        para = tf.add_paragraph()
        para.text = item
        para.level = 0
        for run in para.runs:
            run.font.size = Pt(font_size)
            if font_name:
                run.font.name = font_name
    if highlight is not None:
        for para in tf.paragraphs:
            highlight_paragraph(para, highlight)


def _find_placeholder(slide, *types):
    for ph in slide.placeholders:
        try:
            ph_type = str(ph.placeholder_format.type).split(".")[-1].lower()
        except Exception:
            ph_type = ""
        if any(t in ph_type for t in types):
            return ph
    return None


def _title_placeholder(slide):
    return _find_placeholder(slide, "title")


def _body_placeholders(slide):
    result = []
    for ph in slide.placeholders:
        try:
            ph_type = str(ph.placeholder_format.type).split(".")[-1].lower()
        except Exception:
            ph_type = ""
        if "title" in ph_type or ph_type in {"date", "footer", "slide_number"}:
            continue
        result.append(ph)
    result.sort(key=lambda s: (s.top or 0, s.left or 0))
    return result


def _picture_placeholder(slide):
    return _find_placeholder(slide, "picture")


def _fill_cover(slide, spec, meta_subtitle, bundle=None, *, font_name=None):
    if (t := _title_placeholder(slide)) is not None:
        _set_text(t, spec.title, font_name=font_name)
    sub = _find_placeholder(slide, "subtitle")
    if sub is not None:
        _set_text(sub, spec.subtitle or meta_subtitle or "", font_name=font_name)
    if bundle is None or bundle.meta is None:
        return
    for ph in slide.placeholders:
        try:
            ph_type = str(ph.placeholder_format.type).split(".")[-1].lower()
        except Exception:
            continue
        if "date" in ph_type and bundle.meta.date:
            _set_text(ph, bundle.meta.date, font_name=font_name)
        elif "footer" in ph_type and bundle.meta.author:
            _set_text(ph, bundle.meta.author, font_name=font_name)


def _fill_section(slide, spec, *, font_name=None):
    if (t := _title_placeholder(slide)) is not None:
        _set_text(t, spec.title, font_name=font_name)
    bodies = _body_placeholders(slide)
    if bodies and spec.subtitle:
        _set_text(bodies[0], spec.subtitle, font_name=font_name)


def _fill_content(slide, spec, *, highlight=None, font_name=None):
    if (t := _title_placeholder(slide)) is not None:
        _set_text(t, spec.title, font_name=font_name)
    bodies = _body_placeholders(slide)
    if spec.bullets:
        if len(bodies) >= 2 and len(spec.bullets) >= 5:
            half = (len(spec.bullets) + 1) // 2
            _set_bullets(bodies[0], spec.bullets[:half], highlight=highlight, font_name=font_name)
            _set_bullets(bodies[1], spec.bullets[half:], highlight=highlight, font_name=font_name)
        elif bodies:
            _set_bullets(bodies[0], spec.bullets, highlight=highlight, font_name=font_name)
        else:
            box = SmartLayout(slide).content_region()
            _add_bullet_box(slide, box, spec.bullets, font_name=font_name)
    elif spec.body_text:
        if bodies:
            _set_text(bodies[0], spec.body_text, highlight=highlight, font_name=font_name)
        else:
            box = SmartLayout(slide).content_region()
            _add_text_box(slide, box, spec.body_text, font_size=18, font_name=font_name)


def _fill_two_col(slide, spec, *, highlight=None, font_name=None):
    if (t := _title_placeholder(slide)) is not None:
        _set_text(t, spec.title, font_name=font_name)
    bodies = _body_placeholders(slide)
    bullets = list(spec.bullets)
    half = max(1, (len(bullets) + 1) // 2)
    left, right = bullets[:half], bullets[half:]
    if bodies:
        _set_bullets(bodies[0], left, highlight=highlight, font_name=font_name)
    if len(bodies) > 1:
        _set_bullets(bodies[1], right or [spec.body_text or ""], highlight=highlight, font_name=font_name)
    elif right and not bodies:
        layout = SmartLayout(slide)
        lbox, rbox = layout.two_col_regions()
        _add_bullet_box(slide, lbox, left, font_name=font_name)
        _add_bullet_box(slide, rbox, right, font_name=font_name)
    elif right:
        layout = SmartLayout(slide)
        _, rbox = layout.two_col_regions()
        _add_bullet_box(slide, rbox, right, font_name=font_name)


def _fill_image(slide, spec, bundle, base_dir, *, highlight=None, font_name=None):
    if (t := _title_placeholder(slide)) is not None:
        _set_text(t, spec.title, font_name=font_name)

    img_path = None
    if spec.image_refs:
        asset = bundle.image_by_id(spec.image_refs[0])
        if asset is not None:
            img_path = resolve_image_path(asset, base_dir=base_dir)

    layout = SmartLayout(slide)
    pic_ph = _picture_placeholder(slide)
    has_text = bool(spec.bullets or spec.body_text)

    if pic_ph is not None and img_path is not None:
        from .image_utils import fit_image_for_box

        fitted = fit_image_for_box(
            img_path,
            target_width_emu=int(pic_ph.width or 0),
            target_height_emu=int(pic_ph.height or 0),
        )
        try:
            pic_ph.insert_picture(fitted)
        except Exception:
            slide.shapes.add_picture(
                fitted, pic_ph.left, pic_ph.top, pic_ph.width, pic_ph.height
            )
    elif img_path is not None:
        from .image_utils import fit_image_for_box

        img_region = layout.image_region(has_side_text=has_text)
        fitted = fit_image_for_box(
            img_path,
            target_width_emu=img_region.width,
            target_height_emu=img_region.height,
        )
        slide.shapes.add_picture(
            fitted, img_region.left, img_region.top,
            img_region.width, img_region.height,
        )

    bodies = _body_placeholders(slide)
    text_targets = [b for b in bodies if b is not pic_ph]
    if text_targets:
        if spec.bullets:
            _set_bullets(text_targets[0], spec.bullets, highlight=highlight, font_name=font_name)
        elif spec.body_text:
            _set_text(text_targets[0], spec.body_text, highlight=highlight, font_name=font_name)
    elif has_text and not pic_ph:
        text_region = layout.side_text_region()
        if spec.bullets:
            _add_bullet_box(slide, text_region, spec.bullets, font_name=font_name)
        elif spec.body_text:
            _add_text_box(slide, text_region, spec.body_text, font_size=16, font_name=font_name)


def _fill_thanks(slide, spec, *, font_name=None):
    if (t := _title_placeholder(slide)) is not None:
        _set_text(t, spec.title or "Thanks", font_name=font_name)
    else:
        box = SmartLayout(slide).content_region()
        box.top = int(box.top + box.height * 0.3)
        box.height = int(box.height * 0.4)
        _add_text_box(slide, box, spec.title or "Thanks", font_size=36, font_name=font_name)
    sub = _find_placeholder(slide, "subtitle")
    if sub is not None and spec.subtitle:
        _set_text(sub, spec.subtitle, font_name=font_name)


@dataclass
class _Box:
    left: int
    top: int
    width: int
    height: int


class SmartLayout:
    def __init__(self, slide):
        from pptx.util import Inches, Emu

        prs = slide.part.package.presentation_part.presentation
        self.sw = int(prs.slide_width or Emu(Inches(13.333)))
        self.sh = int(prs.slide_height or Emu(Inches(7.5)))
        self.margin = int(Inches(0.65))
        self.gap = int(Inches(0.25))

        self.title_ph = _find_placeholder(slide, "title", "subtitle")
        if self.title_ph is not None:
            title_bottom = int((self.title_ph.top or 0) + (self.title_ph.height or 0))
            self.title_bottom = title_bottom
            self.title_present = True
        else:
            self.title_bottom = self.margin
            self.title_present = False

    def content_region(self):
        top = self.title_bottom + int(self.gap * 1.5) if self.title_present else self.margin
        bottom = self.sh - self.margin
        left = self.margin
        right = self.sw - self.margin
        h = max(1, bottom - top)
        w = max(1, right - left)
        return _Box(left, top, w, h)

    def image_region(self, has_side_text=False):
        if has_side_text:
            split = self.sw * 3 // 5
            left = self.margin
            top = self.title_bottom + int(self.gap * 1.5) if self.title_present else self.margin
            w = split - self.margin - self.gap // 2
            h = self.sh - top - self.margin
            return _Box(left, top, max(1, w), max(1, h))
        else:
            top = self.title_bottom + int(self.gap * 1.5) if self.title_present else self.margin
            bottom = self.sh - self.margin
            left = self.margin
            right = self.sw - self.margin
            h = max(1, bottom - top)
            w = max(1, right - left)
            return _Box(left, top, w, h)

    def side_text_region(self):
        split = self.sw * 3 // 5
        left = split + self.gap // 2
        top = self.title_bottom + int(self.gap * 1.5) if self.title_present else self.margin
        w = self.sw - left - self.margin
        h = self.sh - top - self.margin
        return _Box(left, top, max(1, w), max(1, h))

    def chart_region(self):
        return self.content_region()

    def two_col_regions(self):
        cr = self.content_region()
        half = (cr.width - self.gap) // 2
        left_box = _Box(cr.left, cr.top, half, cr.height)
        right_box = _Box(cr.left + half + self.gap, cr.top, cr.width - half - self.gap, cr.height)
        return left_box, right_box


def _add_text_box(slide, box, text, *, font_size=None, font_name=None):
    from pptx.util import Pt as _Pt

    txBox = slide.shapes.add_textbox(box.left, box.top, box.width, box.height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.text = text or ""
    for p in tf.paragraphs:
        for run in p.runs:
            if font_size:
                run.font.size = _Pt(font_size)
            if font_name:
                run.font.name = font_name
    return txBox


def _add_bullet_box(slide, box, bullets, *, max_size=28, min_size=14, font_name=None):
    from pptx.util import Pt as _Pt

    txBox = slide.shapes.add_textbox(box.left, box.top, box.width, box.height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.clear()

    box_height_emu = box.height
    if box_height_emu > 0 and len(bullets) > 1:
        box_height_pt = box_height_emu / 12700
        font_size = max(min_size, min(max_size, box_height_pt / (len(bullets) * 1.5 * 1.15)))
    else:
        font_size = max_size

    for idx, item in enumerate(bullets):
        para = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        para.text = item
        para.level = 0
        for run in para.runs:
            run.font.size = _Pt(font_size)
            if font_name:
                run.font.name = font_name
    return txBox


def _chart_box_for_layout(slide):
    layout = SmartLayout(slide)
    return layout.chart_region()


class SmartDecoration:
    def __init__(self, slide, *, theme_color="#3366CC", accent_color="#FF8C00"):
        from pptx.util import Inches, Emu
        from pptx.dml.color import RGBColor as _RGBColor

        prs = slide.part.package.presentation_part.presentation
        self.sw = int(prs.slide_width or Emu(Inches(13.333)))
        self.sh = int(prs.slide_height or Emu(Inches(7.5)))
        self.slide = slide
        self.theme = theme_color
        self.accent = accent_color
        self._tc = hex_to_rgb(theme_color)
        self._ac = hex_to_rgb(accent_color)
        self._light = _RGBColor(235, 240, 248)
        self._emus = Emu

    def _add_rect(self, left, top, width, height, fill_color, line_color=None, line_width=None):
        from pptx.enum.shapes import MSO_SHAPE
        shape = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        if line_color:
            shape.line.color.rgb = line_color
            if line_width:
                shape.line.width = line_width
        else:
            shape.line.fill.background()
        return shape

    def _add_circle(self, left, top, diameter, fill_color, line_color=None):
        from pptx.enum.shapes import MSO_SHAPE
        shape = self.slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, diameter, diameter)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        if line_color:
            shape.line.color.rgb = line_color
        else:
            shape.line.fill.background()
        return shape

    def _add_rounded_rect(self, left, top, width, height, fill_color):
        from pptx.enum.shapes import MSO_SHAPE
        shape = self.slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        return shape

    def _accent_line_h(self, y, left=None, right=None, thickness_pt=2, color=None):
        from pptx.util import Inches, Pt
        l = left or int(Inches(0.65))
        r = right or self.sw - int(Inches(0.65))
        thick = Pt(thickness_pt)
        self._add_rect(l, y - thick // 2, r - l, thick, color or self._ac)

    def _accent_line_v(self, x, top, bottom, thickness_pt=2, color=None):
        from pptx.util import Pt
        thick = Pt(thickness_pt)
        self._add_rect(x - thick // 2, top, thick, bottom - top, color or self._ac)

    def _title_underline(self, title_ph):
        from pptx.util import Inches
        if title_ph is None:
            return
        y = int(title_ph.top or 0) + int(title_ph.height or 0) + int(Inches(0.08))
        left = int(title_ph.left or Inches(0.65))
        width = min(int(title_ph.width or Inches(9)), int(Inches(3.5)))
        self._accent_line_h(y, left=left, right=left + width, thickness_pt=3)

    def _corner_dot(self, x, y, radius_pt=4, color=None):
        from pptx.util import Pt
        d = Pt(radius_pt * 2)
        self._add_circle(x - Pt(radius_pt), y - Pt(radius_pt), d, color or self._tc)

    def decorate_cover(self):
        from pptx.util import Inches
        bar_h = int(Inches(0.12))
        self._add_rect(0, self.sh - bar_h, self.sw, bar_h, self._tc)
        bar_h2 = int(Inches(0.04))
        self._add_rect(0, self.sh - bar_h - bar_h2 - int(Inches(0.02)), self.sw, bar_h2, self._ac)
        import random
        rng = random.Random(42)
        for _ in range(5):
            cx = int(Inches(rng.uniform(0.3, 12.5)))
            cy = int(Inches(rng.uniform(0.3, 6.5)))
            r = int(Inches(rng.uniform(0.06, 0.18)))
            self._add_circle(cx, cy, r, self._ac)

    def decorate_section(self):
        from pptx.util import Inches
        bar_w = int(Inches(0.08))
        self._add_rect(0, 0, bar_w, self.sh, self._tc)
        accent_h = int(Inches(0.03))
        self._add_rect(bar_w, self.sh - int(Inches(0.5)), self.sw - bar_w, accent_h, self._ac)

    def decorate_content(self, title_ph=None):
        if title_ph:
            self._title_underline(title_ph)

    def decorate_toc(self):
        from pptx.util import Inches
        bar_h = int(Inches(0.04))
        self._add_rect(0, self.sh - bar_h, self.sw, bar_h, self._ac)

    def decorate_chart_page(self):
        from pptx.util import Inches
        bar_h = int(Inches(0.04))
        self._add_rect(int(Inches(0.65)), int(Inches(0.35)), int(Inches(1.2)), bar_h, self._tc)

    def decorate_image_page(self):
        from pptx.util import Inches
        margin = int(Inches(0.65))
        dot_r = int(Inches(0.06))
        self._add_circle(margin, margin, dot_r, self._ac)
        self._add_circle(self.sw - margin - dot_r, margin, dot_r, self._tc)

    def decorate_thanks(self):
        from pptx.util import Inches
        import random
        bar_h = int(Inches(0.12))
        self._add_rect(0, self.sh - bar_h, self.sw, bar_h, self._tc)
        rng = random.Random(99)
        for _ in range(8):
            cx = int(Inches(rng.uniform(0.5, 12.5)))
            cy = int(Inches(rng.uniform(0.5, 6.0)))
            r = int(Inches(rng.uniform(0.05, 0.15)))
            color = self._tc if rng.random() > 0.5 else self._ac
            self._add_circle(cx, cy, r, color)


def _style_chart(chart, *, theme_color, accent_color):
    palette = [theme_color, accent_color, "#888888", "#BFBFBF", "#5B9BD5", "#ED7D31"]
    series_list = list(chart.series)
    if len(series_list) == 1:
        only = series_list[0]
        try:
            for i, point in enumerate(only.points):
                fill = point.format.fill
                fill.solid()
                fill.fore_color.rgb = hex_to_rgb(palette[i % len(palette)])
        except Exception:
            try:
                fill = only.format.fill
                fill.solid()
                fill.fore_color.rgb = hex_to_rgb(palette[0])
            except Exception:
                pass
    else:
        for i, series in enumerate(series_list):
            color = palette[i % len(palette)]
            try:
                fill = series.format.fill
                fill.solid()
                fill.fore_color.rgb = hex_to_rgb(color)
            except Exception:
                pass
            try:
                series.format.line.color.rgb = hex_to_rgb(color)
            except Exception:
                pass
    try:
        chart.has_title = False
    except Exception:
        pass


def _fill_chart(slide, spec, *, theme_color, accent_color, font_name=None):
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    if (t := _title_placeholder(slide)) is not None:
        _set_text(t, spec.title, font_name=font_name)

    if spec.chart is None or not spec.chart.categories or not spec.chart.series:
        bodies = _body_placeholders(slide)
        if bodies and spec.bullets:
            _set_bullets(bodies[0], spec.bullets, font_name=font_name)
        elif spec.bullets:
            box = _chart_box_for_layout(slide)
            _add_bullet_box(slide, box, spec.bullets, font_name=font_name)
        return

    box = _chart_box_for_layout(slide)

    cd = CategoryChartData()
    cd.categories = spec.chart.categories
    for ser in spec.chart.series:
        name = str(ser.get("name", ""))
        values = list(ser.get("values", []))
        cd.add_series(name, values)

    kind_map = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "pie": XL_CHART_TYPE.PIE,
    }
    chart_type = kind_map.get(spec.chart.kind, XL_CHART_TYPE.COLUMN_CLUSTERED)

    chart_shape = slide.shapes.add_chart(
        chart_type, box.left, box.top, box.width, box.height, cd
    )
    _style_chart(chart_shape.chart, theme_color=theme_color, accent_color=accent_color)


def _fill_toc(slide, spec, *, highlight=None, font_name=None):
    if (t := _title_placeholder(slide)) is not None:
        _set_text(t, spec.title or "目录", font_name=font_name)
    bodies = _body_placeholders(slide)
    if not spec.bullets:
        return
    if bodies:
        tf = bodies[0].text_frame
        tf.clear()
        for i, item in enumerate(spec.bullets, start=1):
            para = tf.paragraphs[0] if i == 1 else tf.add_paragraph()
            para.text = item
            para.level = 0
            for run in para.runs:
                run.font.size = Pt(18)
                if font_name:
                    run.font.name = font_name
        if highlight is not None:
            for para in tf.paragraphs:
                highlight_paragraph(para, highlight)
    else:
        box = SmartLayout(slide).content_region()
        _add_bullet_box(slide, box, spec.bullets, max_size=20, min_size=16, font_name=font_name)


def _normalize_spacing(slide, line_spacing=1.3, space_after_pt=6):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                p.line_spacing = line_spacing
                p.space_after = Pt(space_after_pt)


def render_to_pptx(
    specs,
    out_path,
    *,
    template_path,
    bundle=None,
    layout_map_yaml=None,
    base_dir=None,
    strip_existing_slides=True,
    theme_color="#3366CC",
    accent_color="#FF8C00",
    enable_polish=True,
    font_name=None,
    enable_decoration=False,
    enable_ai_decoration=False,
):
    if enable_ai_decoration:
        from .ai_decorator import generate_icons_for_outline

        ai_icons = generate_icons_for_outline(list(specs), theme_color=theme_color)
    else:
        ai_icons = {}

    prs = load_template(template_path)
    if strip_existing_slides:
        _strip_existing_slides(prs)
    layout_map = resolve_layout_map(prs, layout_map_yaml)
    bundle = bundle or InputBundle()
    base_dir = base_dir or Path.cwd()

    meta_subtitle = None
    if bundle.meta:
        bits = [bundle.meta.author, bundle.meta.lab, bundle.meta.date]
        meta_subtitle = " · ".join([b for b in bits if b])

    highlight = (
        HighlightStyle(bold=True, color_hex=theme_color) if enable_polish else None
    )

    for page_num, spec in enumerate(specs, start=1):
        layout = layout_map.get(spec.layout)
        slide = prs.slides.add_slide(layout)

        if spec.layout == "cover":
            _fill_cover(slide, spec, meta_subtitle, bundle=bundle, font_name=font_name)
        elif spec.layout == "toc":
            _fill_toc(slide, spec, highlight=highlight, font_name=font_name)
        elif spec.layout == "section":
            _fill_section(slide, spec, font_name=font_name)
        elif spec.layout == "two_col":
            _fill_two_col(slide, spec, highlight=highlight, font_name=font_name)
        elif spec.layout in ("image", "image_only"):
            _fill_image(slide, spec, bundle, base_dir, highlight=highlight, font_name=font_name)
        elif spec.layout == "thanks":
            _fill_thanks(slide, spec, font_name=font_name)
        elif spec.layout == "chart":
            _fill_chart(slide, spec, theme_color=theme_color, accent_color=accent_color, font_name=font_name)
        else:
            _fill_content(slide, spec, highlight=highlight, font_name=font_name)

        if enable_decoration:
            deco = SmartDecoration(slide, theme_color=theme_color, accent_color=accent_color)
            title_ph = _title_placeholder(slide)
            if spec.layout == "cover":
                deco.decorate_cover()
            elif spec.layout == "toc":
                deco.decorate_toc()
            elif spec.layout == "section":
                deco.decorate_section()
            elif spec.layout == "thanks":
                deco.decorate_thanks()
            elif spec.layout == "chart":
                deco.decorate_chart_page()
            elif spec.layout in ("image", "image_only"):
                deco.decorate_image_page()
            else:
                deco.decorate_content(title_ph=title_ph)

        if ai_icons and spec.title in ai_icons and spec.layout == "section":
            from pptx.util import Inches as _In

            icon_w = int(_In(1.4))
            icon_h = int(_In(1.4))
            sl = SmartLayout(slide)
            icon_left = sl.sw - sl.margin - icon_w
            icon_top = int(sl.sh * 0.15)
            slide.shapes.add_picture(
                ai_icons[spec.title], icon_left, icon_top, icon_w, icon_h
            )

        if spec.notes:
            slide.notes_slide.notes_text_frame.text = spec.notes

        if enable_polish and spec.layout not in ("cover", "thanks"):
            fill_footer(
                slide,
                author=bundle.meta.author if bundle.meta else None,
                date=bundle.meta.date if bundle.meta else None,
                page_num=page_num,
            )

        _normalize_spacing(slide)

    out = Path(out_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return str(out)
