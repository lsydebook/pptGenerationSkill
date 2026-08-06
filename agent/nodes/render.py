"""render: draw PPTX from WeeklySlideSpec list using LayoutEngine + python-pptx."""
from __future__ import annotations

import logging
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

from tools.layout_engine import LayoutEngine, get_engine
from tools.background_engine import get_background

from ..state import AgentState

logger = logging.getLogger(__name__)

SLIDE_W = 12192000
SLIDE_H = 6858000
MARGIN_L = 900000
MARGIN_R = 900000
MARGIN_B = 540000
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R

C_CALLOUT_BG = "E8F0F8"
C_CALLOUT_BAR = "3A7BBF"
C_CALLOUT_LABEL = "1A1A2E"
C_CALLOUT_TEXT = "333333"
C_TABLE_HEADER_BG = "1A1A2E"
C_TABLE_HEADER_FG = "FFFFFF"
C_TABLE_ROW_ODD = "F5F7FA"
C_TABLE_ROW_EVEN = "FFFFFF"
C_TABLE_BORDER = "D0D5DD"

_ROLE_MAP = {
    "cover": "cover",
    "thanks": "thanks",
    "section": "section",
}


def _slide_role(layout_name: str) -> str:
    return _ROLE_MAP.get(layout_name, "content")


def _add_rich_background(slide, layout_name: str, accent: str | None = None):
    role = _slide_role(layout_name)
    bg = get_background(role, accent)
    for s in bg.shapes:
        st = 2 if s.shape_type == "oval" else 1
        shape = slide.shapes.add_shape(st, Emu(s.x), Emu(s.y), Emu(s.w), Emu(s.h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb(s.color)
        shape.line.fill.background()


def _hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _add_decorations(slide, decorations):
    for d in decorations:
        st = 2 if d.shape_type == "oval" else 1
        shape = slide.shapes.add_shape(st, Emu(d.x), Emu(d.y), Emu(d.w), Emu(d.h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb(d.color)
        shape.line.fill.background()


def _add_textbox(slide, box, prs):
    from pptx.util import Emu
    if box.bg_color:
        bg_shape = slide.shapes.add_shape(
            1, Emu(box.x), Emu(box.y), Emu(box.w), Emu(box.h),
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = _hex_to_rgb(box.bg_color)
        bg_shape.line.fill.background()
    txBox = slide.shapes.add_textbox(Emu(box.x), Emu(box.y), Emu(box.w), Emu(box.h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = box.text
    p.font.size = Pt(box.font_size)
    p.font.bold = box.bold
    p.font.color.rgb = _hex_to_rgb(box.color)
    p.font.name = box.font_name
    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    p.alignment = align_map.get(box.alignment, PP_ALIGN.LEFT)
    return txBox


def _add_image(slide, box):
    try:
        slide.shapes.add_picture(box.image_path, Emu(box.x), Emu(box.y), Emu(box.w), Emu(box.h))
    except Exception as e:
        logger.warning("failed to add image %s: %s", box.image_path, e)


def _add_chart_shape(slide, spec, has_table=False):
    if not spec.chart or not spec.chart.categories:
        return
    chart = spec.chart
    try:
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE

        chart_data = CategoryChartData()
        chart_data.categories = chart.categories
        for ser in chart.series:
            chart_data.add_series(ser.get("name", ""), ser.get("values", []))

        chart_types = {
            "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line": XL_CHART_TYPE.LINE_MARKERS,
            "pie": XL_CHART_TYPE.PIE,
        }
        chart_type = chart_types.get(chart.kind, XL_CHART_TYPE.COLUMN_CLUSTERED)

        chart_x, chart_y, chart_w, chart_h = 900000, 1400000, 10400000, 4500000
        if has_table:
            chart_h = 3200000

        frame = slide.shapes.add_chart(
            chart_type, Emu(chart_x), Emu(chart_y), Emu(chart_w), Emu(chart_h),
            chart_data,
        )
        chart_obj = frame.chart
        if chart.y_axis_title:
            chart_obj.value_axis.axis_title.text_frame.paragraphs[0].text = chart.y_axis_title
    except Exception as e:
        logger.warning("native chart failed: %s", e)


def _add_callout_box(slide, spec):
    if not spec.callout or not spec.callout.text:
        return

    box_x = Emu(MARGIN_L)
    box_y = Emu(SLIDE_H - MARGIN_B - 380000)
    box_w = Emu(CONTENT_W)
    box_h = Emu(320000)
    bar_w = Emu(80000)

    shape = slide.shapes.add_shape(1, box_x, box_y, box_w, box_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(C_CALLOUT_BG)
    shape.line.fill.background()

    bar = slide.shapes.add_shape(1, box_x, box_y, bar_w, box_h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = _hex_to_rgb(C_CALLOUT_BAR)
    bar.line.fill.background()

    txText = slide.shapes.add_textbox(
        box_x + bar_w + Emu(120000), box_y + Emu(40000),
        box_w - bar_w - Emu(200000), box_h - Emu(80000),
    )
    p = txText.text_frame.paragraphs[0]
    p.text = spec.callout.text
    p.font.size = Pt(13)
    p.font.color.rgb = _hex_to_rgb(C_CALLOUT_TEXT)
    p.font.name = "Microsoft YaHei"


def _add_table_shape(slide, spec):
    if not spec.table or not spec.table.headers or not spec.table.rows:
        return

    tbl = spec.table
    n_rows = len(tbl.rows) + 1
    n_cols = len(tbl.headers)

    table_w = Emu(CONTENT_W)
    table_h = Emu(380000 * n_rows + 100000)

    has_chart = spec.chart is not None
    has_callout = spec.callout is not None

    if has_chart and not has_callout:
        table_y = Emu(4850000)
    elif has_chart and has_callout:
        table_y = Emu(SLIDE_H - MARGIN_B - 500000 - table_h + 100000)
    elif has_callout:
        table_y = Emu(SLIDE_H - MARGIN_B - 700000 - table_h)
    else:
        table_y = Emu(SLIDE_H - MARGIN_B - 500000 - table_h + 100000)

    shape = slide.shapes.add_table(
        n_rows, n_cols, Emu(MARGIN_L), table_y, table_w, table_h
    )
    table_obj = shape.table
    col_w = CONTENT_W // n_cols
    for ci in range(n_cols):
        table_obj.columns[ci].width = Emu(col_w)

    for ci, h in enumerate(tbl.headers):
        cell = table_obj.cell(0, ci)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = _hex_to_rgb(C_TABLE_HEADER_BG)
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = _hex_to_rgb(C_TABLE_HEADER_FG)
            p.font.name = "Microsoft YaHei"
            p.alignment = PP_ALIGN.CENTER

    for ri, row in enumerate(tbl.rows):
        for ci, val in enumerate(row):
            cell = table_obj.cell(ri + 1, ci)
            cell.text = str(val)
            bg = C_TABLE_ROW_ODD if ri % 2 == 0 else C_TABLE_ROW_EVEN
            cell.fill.solid()
            cell.fill.fore_color.rgb = _hex_to_rgb(bg)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = _hex_to_rgb("333333")
                p.font.name = "Microsoft YaHei"
                p.alignment = PP_ALIGN.CENTER


def render(state: AgentState) -> dict:
    specs = state.get("slide_specs") or []
    output_path = state.get("output_path", "out/weekly.pptx")
    accent = state.get("accent_color", "")

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)

    engine = get_engine(accent if accent else None)

    for spec in specs:
        layout = engine.compute(spec)

        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        bg = layout.bg_color
        if bg and bg.upper() != "FFFFFF":
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = _hex_to_rgb(bg)

        _add_rich_background(slide, spec.layout, accent)

        for tb in layout.text_boxes:
            _add_textbox(slide, tb, prs)

        for ib in layout.image_boxes:
            _add_image(slide, ib)

        if layout.chart_area and spec.chart:
            _add_chart_shape(slide, spec, has_table=spec.table is not None)

        if spec.table:
            _add_table_shape(slide, spec)

        if spec.callout:
            _add_callout_box(slide, spec)

    prs.save(str(out))
    logger.info("saved: %s (%d slides)", out, len(specs))
    return {"rendered_path": str(out)}
